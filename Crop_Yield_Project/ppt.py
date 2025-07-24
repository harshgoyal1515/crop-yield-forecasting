from pptx import Presentation
from pptx.util import Inches
import os

# Create a presentation object
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Intelligent Crop Yield Forecasting"
subtitle.text = "Using Big Data Analytics\nIIT Jodhpur - Big Data Management Project"

# Team slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Team Members"
slide.placeholders[1].text = "• Abhinav Tote\n• Harsh Goyal\n• Harshal Ingle\n• Omar Khan"

# Phase 1 - EDA
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Phase 1: Exploratory Data Analysis (EDA)"
slide.placeholders[1].text = (
    "• Analyzed historical crop and rainfall data\n"
    "• Visualized yield distributions for 23 major crops\n"
    "• Cleaned and standardized inconsistent column names\n"
    "• Plots saved in: `eda_plots/`"
)

# Insert a sample EDA image
eda_image = "eda_plots/yield_distribution.png"
if os.path.exists(eda_image):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Sample EDA Plot"
    slide.shapes.add_picture(eda_image, Inches(1), Inches(1.5), width=Inches(7.5))

# Phase 2 - Big Data with PySpark
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Phase 2: Big Data Processing with Spark"
slide.placeholders[1].text = (
    "• Set up PySpark on local system\n"
    "• Successfully created SparkSession\n"
    "• Explored yield trends using Spark transformations\n"
    "• Filtered and grouped large data efficiently"
)

# Phase 3 - Machine Learning
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Phase 3: ML Prediction"
slide.placeholders[1].text = (
    "• Used Linear Regression to predict rice yield\n"
    "• Features: area (1000 ha) and monsoon rainfall (Jun-Sep)\n"
    "• Achieved decent R² score after data cleaning\n"
    "• Focused on rice as case study crop"
)

# Sample prediction result
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Prediction Results"
slide.placeholders[1].text = (
    "• Model: Linear Regression\n"
    "• Target: RICE YIELD (Kg per ha)\n"
    "• Input: ['JUN', 'JUL', 'AUG', 'SEP', 'RICE AREA (1000 ha)']\n"
    "• Train-test split: 80-20\n"
    "• R² Score and actual vs predicted plots analyzed"
)

# Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion & Future Work"
slide.placeholders[1].text = (
    "• Framework for crop yield forecasting is established\n"
    "• Further work can include:\n"
    "    - Adding more crops & models\n"
    "    - Integration with live weather APIs\n"
    "    - Deploying model using MLOps"
)

# Save the presentation
prs.save("Crop_Yield_Project_Presentation.pptx")
print("✅ Presentation saved as Crop_Yield_Project_Presentation.pptx")
